"""Document Ingestion Engine"""
import os
import uuid
from typing import Dict, Any, Optional
from pathlib import Path
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
from app.models.document import DocumentType


class IngestionEngine:
    """Document Ingestion Engine for multi-format parsing"""
    
    def __init__(self):
        self.supported_formats = {
            DocumentType.PDF: self._parse_pdf,
            DocumentType.DOCX: self._parse_docx,
            DocumentType.PPTX: self._parse_pptx,
            DocumentType.TXT: self._parse_txt,
            DocumentType.MD: self._parse_txt,
            DocumentType.HTML: self._parse_html,
        }
    
    def detect_document_type(self, filename: str) -> DocumentType:
        """Detect document type from filename"""
        ext = filename.rsplit(".", 1)[1].lower() if "." in filename else ""
        ext_map = {
            "pdf": DocumentType.PDF,
            "docx": DocumentType.DOCX,
            "pptx": DocumentType.PPTX,
            "txt": DocumentType.TXT,
            "md": DocumentType.MD,
            "html": DocumentType.HTML,
            "htm": DocumentType.HTML,
        }
        return ext_map.get(ext, DocumentType.TXT)
    
    def ingest_document(
        self,
        file_path: str,
        filename: str,
        user_id: str,
        title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Ingest a document and extract its content
        
        Args:
            file_path: Path to the document file
            filename: Original filename
            user_id: User ID who uploaded the document
            title: Optional title (defaults to filename)
        
        Returns:
            Dictionary containing document metadata and content
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        doc_type = self.detect_document_type(filename)
        file_size = os.path.getsize(file_path)
        
        # Parse document content
        parser = self.supported_formats.get(doc_type)
        if not parser:
            raise ValueError(f"Unsupported document type: {doc_type}")
        
        content = parser(file_path)
        
        # Analyze document
        analysis = self._analyze_document(content, doc_type)
        
        # Generate document ID
        doc_id = str(uuid.uuid4())
        
        return {
            "id": doc_id,
            "user_id": user_id,
            "title": title or filename,
            "filename": filename,
            "file_path": file_path,
            "file_size": file_size,
            "document_type": doc_type,
            "content": content,
            "metadata": {
                "page_count": analysis.get("page_count", 0),
                "word_count": analysis.get("word_count", 0),
                "char_count": analysis.get("char_count", 0),
                "has_tables": analysis.get("has_tables", False),
                "has_code": analysis.get("has_code", False),
                "language": analysis.get("language", "en"),
            },
            "domain": analysis.get("domain", "general"),
            "complexity_score": analysis.get("complexity_score", 1),
            "language": analysis.get("language", "en"),
        }
    
    def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF document"""
        doc = fitz.open(file_path)
        text_content = []
        
        for page in doc:
            text = page.get_text()
            text_content.append(text)
        
        doc.close()
        return "\n\n".join(text_content)
    
    def _parse_docx(self, file_path: str) -> str:
        """Parse DOCX document"""
        doc = DocxDocument(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                text_content.append(row_text)
        
        return "\n\n".join(text_content)
    
    def _parse_pptx(self, file_path: str) -> str:
        """Parse PPTX document"""
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_content.append("\n".join(slide_text))
        
        return "\n\n--- Slide ---\n\n".join(text_content)
    
    def _parse_txt(self, file_path: str) -> str:
        """Parse plain text document"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    
    def _parse_html(self, file_path: str) -> str:
        """Parse HTML document"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html = f.read()
        
        soup = BeautifulSoup(html, "html.parser")
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)
        
        return text
    
    def _analyze_document(self, content: str, doc_type: DocumentType) -> Dict[str, Any]:
        """Analyze document content for metadata"""
        analysis = {
            "page_count": 0,
            "word_count": len(content.split()),
            "char_count": len(content),
            "has_tables": False,
            "has_code": False,
            "language": "en",
            "domain": "general",
            "complexity_score": 1,
        }
        
        # Detect tables
        analysis["has_tables"] = "|" in content or "table" in content.lower()
        
        # Detect code
        code_indicators = ["```", "def ", "function ", "class ", "import ", "#include"]
        analysis["has_code"] = any(indicator in content for indicator in code_indicators)
        
        # Detect domain (simple heuristic)
        legal_keywords = ["contract", "law", "legal", "statute", "regulation"]
        technical_keywords = ["api", "function", "code", "programming", "algorithm"]
        creative_keywords = ["story", "novel", "poem", "creative", "fiction"]
        
        content_lower = content.lower()
        if any(kw in content_lower for kw in legal_keywords):
            analysis["domain"] = "legal"
        elif any(kw in content_lower for kw in technical_keywords):
            analysis["domain"] = "technical"
        elif any(kw in content_lower for kw in creative_keywords):
            analysis["domain"] = "creative"
        
        # Calculate complexity score (1-5)
        avg_word_length = sum(len(word) for word in content.split()) / max(len(content.split()), 1)
        sentence_count = content.count(".") + content.count("!") + content.count("?")
        avg_sentence_length = len(content.split()) / max(sentence_count, 1)
        
        complexity = 1
        if avg_word_length > 5:
            complexity += 1
        if avg_sentence_length > 20:
            complexity += 1
        if analysis["has_code"]:
            complexity += 1
        if analysis["domain"] == "legal" or analysis["domain"] == "technical":
            complexity += 1
        
        analysis["complexity_score"] = min(complexity, 5)
        
        return analysis
